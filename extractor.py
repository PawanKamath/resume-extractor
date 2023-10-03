import csv
import json
import logging
import os
import nltk
from collections import defaultdict

import markdown
import PyPDF2
import requests
import speech_recognition as sr
import whisper
from bs4 import BeautifulSoup
from docx import Document
from moviepy.editor import VideoFileClip
from openpyxl import load_workbook
from paddleocr import PaddleOCR
from pptx import Presentation
from pypdf import PdfReader
from tqdm import tqdm
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
import string
import docx2txt
import shutil

os.environ["PADDLEOCR_DEBUG"] = "False"
os.environ["KMP_WARNINGS"] = "off"  # Suppresses PaddleOCR debug logs
logging.getLogger("ppocr").setLevel(logging.ERROR)

class TextExtractor:
    """
    Initializes the TextExtractor class with a folder path and output CSV file path.
    """

    def __init__(self, folder_path, output_json):
        self.folder_path = folder_path
        self.output_json = output_json
        self.ocr = PaddleOCR()
        logging.basicConfig(
            filename="file_extraction.log",
            level=logging.ERROR,
            format="%(asctime)s: %(levelname)s: %(message)s",
            datefmt="%Y-%m-%d %H:%M:%S",
        )

    def clean_text(self, text: str) -> str:
        translation_table = str.maketrans(string.punctuation, ' ' * len(string.punctuation))
        text = text.translate(translation_table)
        words = word_tokenize(text)
        filler_words = set(stopwords.words('english'))
        words = [word for word in words if word.lower() not in filler_words]

        text = ' '.join(words)
        return text

    def extract_keywords(self, text, num_keywords=5):  
        words = word_tokenize(text)  
        stop_words = set(stopwords.words('english'))  
        filtered_words = [word for word in words if word.casefold() not in stop_words]  
        freq_dist = nltk.FreqDist(filtered_words)  
        keywords = [word for word, _ in freq_dist.most_common(num_keywords)]  
        return keywords 

    def getinfo(self, file_path, text):
        fname = os.path.basename(file_path)
        try:
            if text!="":
                data = self.clean_text(text)
                keywords = self.extract_keywords(data)
                metadata = {"source":file_path,
                            "file_name": fname,
                            "keywords": keywords}
                data_dict = { str(fname): {"metadata":metadata, "data":data} }
                return data_dict
            else:
                return dict()
        except Exception as e:
            logging.info(f"No text {file_path}: {e}")
            return dict()

    def extract_text_from_images(self, file_path):
        """
        Extracts text from the image at the given file path using PaddleOCR.
        Returns the extracted text as a string.
        """
        try:
            result = self.ocr.ocr(file_path, cls=True)
            text = ""
            for line in result:
                if len(line) >= 1 and isinstance(line[0][1], str):
                    text += line[0][1] + "\n"
            return self.getinfo(file_path, text)
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_docx(self, file_path):
        """
        Extracts text from the .docx file at the given file path.
        Returns the extracted text as a string.
        """
        try:
            # doc = Document(file_path)
            # text = "\n".join([para.text for para in doc.paragraphs])
            text = docx2txt.process(file_path)
            return self.getinfo(file_path, text)
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_audio(self, file_path):
        """
        Extracts text from the .mp3, .wav file at the given file path.
        Returns the extracted text as a string.
        """
        try:
            model = whisper.load_model("base")
            response = model.transcribe(file_path)
            return self.getinfo(file_path, response["text"])
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_xlsx(self, file_path):
        """
        Extracts text from the .xlsx file at the given file path.
        Returns the extracted text as a string.
        """
        try:
            wb = load_workbook(file_path)
            text = ""
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                for row in ws.iter_rows():
                    row_text = " ".join([str(cell.value) if cell.value is not None else "" for cell in row])
                    text += row_text + "\n"
            return self.getinfo(file_path, text)
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_pdf(self, file_path):
        """
        Extracts text from the .pdf file at the given file path.
        Returns the extracted text as a string.
        """
        try:
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page_num in range(len(pdf_reader.pages)):
                    text += pdf_reader.pages[page_num].extract_text()
            return self.getinfo(file_path, text)
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_txt(self, file_path):
        """
        Extracts text from the .txt file at the given file path.
        Returns the extracted text as a string.
        """
        try:
            with open(file_path) as file:
                text = file.read()
            return self.getinfo(file_path, text)
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_csv(self, file_path):
        """
        Extracts text from the .csv file at the given file path.
        Returns the extracted text as a string.
        """
        try:
            with open(file_path) as file:
                reader = csv.reader(file)
                text = "\n".join([",".join(row) for row in reader])
            return self.getinfo(file_path, text)
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_html(self, file_path):
        """
        Extracts text from the .html file at the given file path using BeautifulSoup.
        Returns the extracted text as a string.
        """
        try:
            with open(file_path, encoding="utf-8") as file:
                content = file.read()
            soup = BeautifulSoup(content, "html.parser")
            text = soup.get_text(separator="\n")
            return self.getinfo(file_path, text)
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_url(self, file_path):
        """
        Extracts text from the url using BeautifulSoup.
        Returns the extracted text as a string.
        """
        try:
            response = requests.get(file_path)
            if response.status_code == 200:
                soup = BeautifulSoup(response.content, "html.parser")
                text = soup.get_text(separator="\n")
                return self.getinfo(file_path, text)
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_pptx(self, file_path):
        """
        Extracts text from the .pptx file at the given file path.
        Returns the extracted text as a string.
        """
        try:
            prs = Presentation(file_path)
            text = "\n".join(
                [
                    "\n".join([shape.text for shape in slide.shapes if shape.has_text_frame])
                    for slide in prs.slides
                ]
            )
            return self.getinfo(file_path, text)
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_recursive(self, data, result=""):
        """
        Recursively extracts text from the nested data structure (dict, list, or str).
        Returns the extracted text as a string.
        """
        if isinstance(data, dict):
            for key, value in data.items():
                result = self.extract_text_recursive(value, result)
        elif isinstance(data, list):
            for item in data:
                result = self.extract_text_recursive(item, result)
        elif isinstance(data, str):
            result += " " + data
        return result

    def extract_text_from_json(self, file_path):
        """
        Extracts text from the .json file at the given file path.
        Returns the extracted text as a string.
        """
        try:
            with open(file_path) as file:
                data = json.load(file)
            return self.extract_text_recursive(data).strip()
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_markdown(self, file_path):
        """
        Extracts text from the .md file at the given file path.
        Returns the extracted text as a string.
        """
        with open(file_path, encoding="utf-8") as md_file:
            md_content = md_file.read()
        html_content = markdown.markdown(md_content)
        soup = BeautifulSoup(html_content, "html.parser")
        text = soup.get_text()

        return self.getinfo(file_path, text)

    def extract_text_from_mp4(self, file_path):
        """
        Extracts text from the .mp4 file at the given file path.
        Returns the extracted text as a string.
        """
        try:
            video = VideoFileClip(file_path)
            audio = video.audio
            temp_audio_file = "temp_audio.mp3"
            audio.write_audiofile(temp_audio_file)
            model = whisper.load_model("base")
            response = model.transcribe(file_path)
            os.remove(temp_audio_file)

            return self.getinfo(file_path, response["text"])

        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_pdf_and_embedded_images(self, file_path):
        """
        Extracts text from the .pdf file at the given file path.
        Returns the extracted text as a string and saves images.
        """
        try:
            text = ""
            reader = PdfReader(file_path)
            if not os.path.exists("extract_images_pdf"):
                os.makedirs("extract_images_pdf")
            for page in reader.pages:
                text += page.extract_text()
                for i in page.images:
                    with open(os.path.join("extract_images_pdf", i.name), "wb") as f:
                        f.write(i.data) 
                    try:
                        img_ocr = self.extract_text_from_images(os.path.join("extract_images_pdf", i.name))
                        img_data = img_ocr[str(i.name)]["data"]
                        text+=img_data
                    except Exception as e:
                        logging.info(f"No text in {i.name}: {e}")
                        pass
            shutil.rmtree("extract_images_pdf")
            return self.getinfo(file_path, text)

        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return dict()

    def extract_text_from_single_document(self):
        file_types = {
            ".docx": self.extract_text_from_docx,
            ".xlsx": self.extract_text_from_xlsx,
            ".pdf": self.extract_text_from_pdf_and_embedded_images,
            ".txt": self.extract_text_from_txt,
            ".csv": self.extract_text_from_csv,
            ".pptx": self.extract_text_from_pptx,
            ".html": self.extract_text_from_html,
            ".json": self.extract_text_from_json,
            ".png": self.extract_text_from_images,
            ".jpg": self.extract_text_from_images,
            ".jpeg": self.extract_text_from_images,
            ".bmp": self.extract_text_from_images,
            ".tiff": self.extract_text_from_images,
            ".md": self.extract_text_from_markdown,
            ".mp4": self.extract_text_from_mp4,
            ".mp3": self.extract_text_from_audio,
            ".wav": self.extract_text_from_audio,
        }
        # extracted_texts = defaultdict(str)
        file_extension = os.path.splitext(self.folder_path)[1]
        if file_extension in file_types:
            print(f"Processing:{self.folder_path}")
            try:
                extracted_texts = file_types[file_extension](self.folder_path)
                return extracted_texts
            except Exception as e:
                logging.error(f"Error extracting text from {self.folder_path}: {e}")

    def extract_text_from_documents(self):
        """
        Extracts text from supported file types in the specified folder.
        Returns a list of dictionaries containing file names and extracted text.
        """
        file_types = {
            ".docx": self.extract_text_from_docx,
            ".xlsx": self.extract_text_from_xlsx,
            ".pdf": self.extract_text_from_pdf_and_embedded_images,
            ".txt": self.extract_text_from_txt,
            ".csv": self.extract_text_from_csv,
            ".pptx": self.extract_text_from_pptx,
            ".html": self.extract_text_from_html,
            ".json": self.extract_text_from_json,
            ".png": self.extract_text_from_images,
            ".jpg": self.extract_text_from_images,
            ".jpeg": self.extract_text_from_images,
            ".bmp": self.extract_text_from_images,
            ".tiff": self.extract_text_from_images,
            ".md": self.extract_text_from_markdown,
            ".mp4": self.extract_text_from_mp4,
            ".mp3": self.extract_text_from_audio,
            ".wav": self.extract_text_from_audio,
        }

        extracted_texts = defaultdict(str)

        for root, _, files in os.walk(self.folder_path):
            for file in tqdm(files):
                file_path = os.path.join(root, file)
                if "http" in file_path:
                    extracted_texts[file] = self.extract_text_from_url(file_path)
                else:
                    file_extension = os.path.splitext(file)[1]
                    if file_extension in file_types:
                        print(f"Processing:{file}")
                        try:
                            extracted_texts[file] = file_types[file_extension](file_path)
                        except Exception as e:
                            logging.error(f"Error extracting text from {file_path}: {e}")

        return [{"file_name": key, "text": value} for key, value in extracted_texts.items()]

    """
    def save_to_csv(self, data):
        # Saves the extracted data (list of dictionaries) to a CSV file.
        try:
            with open(self.output_csv, "w", newline="", encoding="utf-8") as file:
                writer = csv.DictWriter(file, fieldnames=["file_name", "text"])
                writer.writeheader()
                writer.writerows(data)
        except Exception as e:
            logging.error(f"Error saving to {self.output_csv}: {e}")
    """

    def save_to_json(self, data):  
        try:
            with open(self.output_json, 'w') as f:  
                json.dump(data, f, indent=4) 
        except Exception as e:
            logging.error(f"Error saving to {self.output_json}: {e}") 
